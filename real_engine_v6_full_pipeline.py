"""Fail-closed, source-grounded presentation generation and verification."""

from __future__ import annotations

from collections import defaultdict
import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from atomic_claim_evidence import AtomicClaimBuilder, NUMBER_PATTERN
from atomic_content_verifier_v6 import AtomicContentVerifierV6
from chart_generator_real import ChartGeneratorReal
from evidence_extractor_real import EvidenceExtractorReal
from evidence_retriever_v6 import EvidenceRetrieverV6
from invisible_metadata_engine import InvisibleMetadataEngine
from openai_narrative import OpenAINarrative
from tamper_detector_v6 import TamperDetectorV6
from tavily_research import TavilyResearch


LOGGER = logging.getLogger(__name__)


class RealEngineV6FullPipeline:
    """Preserve the existing public API while enforcing slide-wide verification."""

    def __init__(self, tavily_key=None, openai_key=None, *, narrator=None, researcher=None, signing_key=None):
        self.extractor = EvidenceExtractorReal()
        self.researcher = researcher or (TavilyResearch(api_key=tavily_key) if tavily_key else None)
        self.narrator = narrator or (OpenAINarrative(api_key=openai_key) if openai_key else None)
        self.atomic_builder = AtomicClaimBuilder()
        self.content_verifier = AtomicContentVerifierV6()
        self.tamper_detector = TamperDetectorV6()
        self.retriever = EvidenceRetrieverV6()
        self.invisible_engine = InvisibleMetadataEngine(signing_key=signing_key)
        self.chart_gen = ChartGeneratorReal()
        self.evidence_store = []
        self.atomic_claims = []
        self.research_results = {}

    def _add_fact_evidence(self, fact: dict) -> None:
        if not fact.get("can_use_in_deck", True):
            return
        location = {
            "page": fact.get("page_number"),
            "cell_range": fact.get("cell_range"),
            "cell": fact.get("cell"),
            "sheet": fact.get("sheet"),
            "paragraph": fact.get("paragraph_number"),
            "url": fact.get("url"),
        }
        location = {key: value for key, value in location.items() if value is not None}
        if not location:
            raise ValueError(f"Source {fact.get('source_file', 'unknown')} does not contain a real evidence location")
        start, end = fact.get("char_start"), fact.get("char_end")
        self.atomic_builder.add_atomic_evidence(
            source_file=fact.get("source_file") or fact.get("url"),
            exact_location=location,
            exact_passage=fact.get("source_text") or fact.get("claim"),
            char_span=(start, end) if start is not None and end is not None else None,
            url=fact.get("url"),
        )
        self.evidence_store.append(self.atomic_builder.evidences[-1])

    def step_2_extract(self, uploaded_paths):
        facts = []
        extractors = {
            ".pdf": self.extractor.extract_pdf_with_pages,
            ".xlsx": self.extractor.extract_excel_with_cells,
            ".docx": self.extractor.extract_docx_with_paragraph,
            ".csv": self.extractor.extract_csv_with_rows,
        }
        for filename in uploaded_paths:
            suffix = Path(filename).suffix.casefold()
            if suffix not in extractors:
                raise ValueError(f"Unsupported source format: {suffix or '(no extension)'}")
            extracted = extractors[suffix](str(filename))
            for fact in extracted:
                self._add_fact_evidence(fact)
                facts.append(fact)
        return facts

    def step_3_research(self, idea):
        if self.researcher is None:
            self.research_results = {}
            return {}
        results = self.researcher.research_deck_gaps(idea)
        for response in results.values():
            for item in response.get("results", [])[:5]:
                if not item.get("url") or not item.get("source_text"):
                    continue
                self._add_fact_evidence({"source_file": item["url"], "url": item["url"], "source_text": item["source_text"], "claim": item.get("claim"), "can_use_in_deck": True})
        self.research_results = results
        return results

    def step_5_generate_atomic_claims(self, idea, facts):
        if self.narrator is None:
            raise RuntimeError("OPENAI_API_KEY is required; fabricated fallback slides are disabled")
        if not self.evidence_store:
            raise ValueError("Upload a source document or enable research before generating a verified deck")

        slides = self.narrator.generate_defensible_deck(idea, facts, self.evidence_store)
        claims = []
        for slide_index, slide in enumerate(slides, start=1):
            if NUMBER_PATTERN.search(str(slide.get("title") or "")):
                raise ValueError(f"Slide {slide_index} has an unaudited number in its title; move numeric claims into sourced bullets")
            citations = list(dict.fromkeys(str(item) for item in slide.get("citations", []) if item))
            for bullet in slide.get("bullets", []):
                statement = str(bullet).strip()
                if not statement:
                    continue
                claim = self.atomic_builder.add_atomic_claim(statement, slide_index, "Source: " + ", ".join(citations))
                self.retriever.link_claim_atomically(claim, self.evidence_store)
                mapped_sources = [evidence.source_file for evidence in self.evidence_store if evidence.evidence_id in claim.evidence_ids]
                actual_sources = list(dict.fromkeys(mapped_sources or citations))
                claim.visible_citation = "Source: " + ", ".join(actual_sources) if actual_sources else ""
                claims.append(claim)

            slide_claims = [claim for claim in claims if claim.slide_number == slide_index]
            slide["citations"] = list(dict.fromkeys(source for claim in slide_claims for source in ([evidence.source_file for evidence in self.evidence_store if evidence.evidence_id in claim.evidence_ids] or citations)))

        if not slides or not claims:
            raise RuntimeError("Deck generation did not produce auditable content")
        self.atomic_claims = claims
        return slides, claims

    @staticmethod
    def _clear_template_slides(presentation) -> None:
        slide_id_list = presentation.slides._sldIdLst
        for slide_id in list(slide_id_list):
            relationship_id = slide_id.rId
            presentation.part.drop_rel(relationship_id)
            slide_id_list.remove(slide_id)

    def step_7_create_deck_with_invisible_layer(self, template_path, slides, atomic_claims, output_path):
        template = Path(template_path) if template_path else None
        presentation = Presentation(str(template)) if template and template.is_file() else Presentation()
        self._clear_template_slides(presentation)
        evidence_lookup = {evidence.evidence_id: evidence for evidence in self.evidence_store}

        for index, data in enumerate(slides, start=1):
            layout_index = 1 if len(presentation.slide_layouts) > 1 else 0
            slide = presentation.slides.add_slide(presentation.slide_layouts[layout_index])
            title = slide.shapes.title
            if title is None:
                title = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), presentation.slide_width - Inches(1.2), Inches(0.8))
            title.text = str(data.get("title") or f"Slide {index}")

            claims = [claim for claim in atomic_claims if claim.slide_number == index]
            if not claims:
                raise ValueError(f"Slide {index} contains no auditable claims")

            body = None
            for shape in slide.shapes:
                if shape != title and getattr(shape, "has_text_frame", False):
                    body = shape
                    break
            if body is None:
                body = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), presentation.slide_width - Inches(1.4), presentation.slide_height - Inches(2.4))
            text_frame = body.text_frame
            text_frame.clear()
            for claim_index, claim in enumerate(claims):
                paragraph = text_frame.paragraphs[0] if claim_index == 0 else text_frame.add_paragraph()
                paragraph.text = claim.statement
                paragraph.level = 0
                paragraph.font.size = Pt(20)
                paragraph.space_after = Pt(14)

            citation_sources = list(dict.fromkeys(data.get("citations") or [source for claim in claims for source in claim.visible_citation.removeprefix("Source: ").split(", ") if source]))
            footer = slide.shapes.add_textbox(Inches(0.45), presentation.slide_height - Inches(0.55), presentation.slide_width - Inches(0.9), Inches(0.32))
            footer.text_frame.text = "Source: " + "; ".join(citation_sources) if citation_sources else "Source: unavailable"
            footer.text_frame.paragraphs[0].font.size = Pt(9)
            self.invisible_engine.add_invisible_metadata_to_slide(slide, claims, evidence_lookup=evidence_lookup)

        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(output))
        return str(output)

    def step_8_qc_v6_full(self, pptx_path):
        presentation = Presentation(str(pptx_path))
        results = {"total_slides": len(presentation.slides), "atomic_content_checks": [], "tamper_checks": [], "can_publish": bool(presentation.slides), "qc_pass": bool(presentation.slides), "issues": [], "adversarial_tests": [], "verified_claim_count": 0, "integrity_signed": bool(self.invisible_engine.signing_key)}
        if not presentation.slides:
            results["issues"].append("The presentation contains no slides")

        expected_by_slide = defaultdict(list)
        for claim in self.atomic_claims:
            expected_by_slide[claim.slide_number].append(claim)
        evidence_lookup = {item.evidence_id: item for item in self.evidence_store}

        for index, slide in enumerate(presentation.slides, start=1):
            metadata = self.invisible_engine.extract_invisible_metadata(slide)
            if not metadata:
                results["issues"].append(f"Slide {index}: audit metadata is missing")
                continue
            if not self.invisible_engine.validate_integrity(metadata):
                results["issues"].append(f"Slide {index}: audit metadata integrity validation failed")
                continue

            stored_claims = metadata.get("citedeck_verification", {}).get("claims", [])
            expected = expected_by_slide[index]
            if {record.get("claim_id") for record in stored_claims} != {claim.claim_id for claim in expected}:
                results["issues"].append(f"Slide {index}: the audit layer does not contain every visible claim")
                continue

            tamper = self.tamper_detector.detect_tampering(slide, expected)
            results["tamper_checks"].append(tamper)
            results["issues"].extend(f"Slide {index}: {issue}" for issue in tamper["issues"])

            for claim in expected:
                if not claim.visible_citation:
                    results["issues"].append(f"Slide {index}, {claim.claim_id}: source citation is missing")
                for span in claim.numeric_spans:
                    evidence = evidence_lookup.get(span.evidence_id)
                    if evidence is None:
                        results["issues"].append(f"Slide {index}, {claim.claim_id}: {span.value} has no independently sourced evidence")
                        continue
                    check = self.content_verifier.verify_atomic_claim(span, evidence)
                    check["claim_id"] = claim.claim_id
                    check["slide"] = index
                    results["atomic_content_checks"].append(check)
                    if not check["verified"]:
                        results["issues"].append(f"Slide {index}, {claim.claim_id}: {span.value} is not supported by {evidence.source_file}")
                    else:
                        span.verification_status = "VERIFIED_EXACT_SPAN"
                results["verified_claim_count"] += 1

        expected_count = sum(len(claim.numeric_spans) for claim in self.atomic_claims)
        if len(results["atomic_content_checks"]) != expected_count:
            results["issues"].append(f"Only {len(results['atomic_content_checks'])} of {expected_count} numeric spans were independently checked")

        results["can_publish"] = bool(presentation.slides) and not results["issues"]
        results["qc_pass"] = results["can_publish"]
        results["verification_status"] = "PASSED - every claim and every number has source-backed evidence" if results["can_publish"] else "FAILED - export blocked until all source and integrity issues are resolved"
        results["trust_badge"] = "CiteDeck Verified: source-grounded, slide-wide integrity checked" if results["can_publish"] else "Verification incomplete - export blocked"
        return results
