import streamlit as st
import os, json
from pathlib import Path

from real_engine_v6_full_pipeline import RealEngineV6FullPipeline
from atomic_claim_evidence import AtomicClaimBuilder

st.set_page_config(page_title="CiteDeck V6 - Full Pipeline + Adversarial Tests Fixed", layout="wide")

st.title("CiteDeck V6 — Full End-to-End Pipeline + Adversarial Tests Fixed")
st.caption("Fixes your 4 critical issues: content verification, tamper detection, production retrieval, full pipeline integration")

tavily_key = st.secrets.get("TAVILY_API_KEY") if "TAVILY_API_KEY" in st.secrets else os.getenv("TAVILY_API_KEY")
openai_key = st.secrets.get("OPENAI_API_KEY") if "OPENAI_API_KEY" in st.secrets else os.getenv("OPENAI_API_KEY")

with st.expander("Your 4 critical issues - V6 fixes", expanded=True):
    st.markdown("""
    **Critical 1: V5 does not verify evidence actually supports number**
    Your test: Claim $25B with evidence passage "Market is actually $10B" -> V5 passed (should fail)
    
    V6 fix `atomic_content_verifier_v6.py`:
    ```python
    def does_passage_contain_number(claimed_value, passage):
        # Checks if $25B actually in passage, not just evidence ID exists
        exact_match = claimed_value in passage
        # Checks for conflicting numbers - if passage has $10B but claim $25B, it's conflict
        if not exact_match and conflicting_numbers:
            return {"supports": False, "adversarial_test_detected": True}
    ```
    Now $25B claim with $10B evidence correctly FAILS
    
    **Critical 2: Visible slide can be tampered with**
    Your test: Hidden $25B, visible changed to $999B -> V5 passed
    
    V6 fix `tamper_detector_v6.py`:
    ```python
    visible_numbers = extract_numbers(visible_slide_text) # $999B
    hidden_numbers = [ns.value for ns in claim.numeric_spans] # $25B
    if set(visible_numbers) != set(hidden_numbers):
        tamper_detected = True
        issues.append("Visible $999B != Hidden $25B - TAMPER DETECTED")
    ```
    Now visible $999B vs hidden $25B correctly FAILS and blocks export
    
    **Important 3: Atomic lookup placeholder logic**
    V5: `if value in passage or source in citation` -> false mapping $10M expenses vs revenue
    
    V6 `evidence_retriever_v6.py`:
    ```python
    score = 0
    if exact value in passage: +50
    keyword overlap claim vs passage: +5 per keyword
    context mismatch revenue vs expenses: -20
    suitable only if score >=40
    ```
    Prefers relevant evidence, avoids false $10M mapping
    
    **Important 4: V5 not integrated into full pipeline**
    V5: app_v5_atomic_invisible.py is demo with manually created evidence
    
    V6 `real_engine_v6_full_pipeline.py`:
    ```
    Upload docs -> Extract with page/cell -> Tavily research -> OpenAI narrative
      -> Atomic claims with numeric spans -> EvidenceRetriever links each number to best evidence
      -> InvisibleMetadataEngine embeds in notes + clean visible citation
      -> QCV6: content verification + tamper detection + atomic checks
      -> Export blocked if fails
    ```
    Full chain wired
    """)

# Adversarial tests
st.subheader("Adversarial Tests - These failed in V5, should fail correctly in V6")

col1, col2 = st.columns(2)

with col1:
    st.markdown("**Test 1: Content mismatch $25B claim vs $10B evidence**")
    if st.button("Run Test 1 - Should FAIL correctly"):
        from atomic_content_verifier_v6 import AtomicContentVerifierV6
        from atomic_claim_evidence import AtomicClaimBuilder
        
        builder = AtomicClaimBuilder()
        eid = builder.add_atomic_evidence(
            source_file="market_report.pdf",
            exact_location={"page": 14},
            exact_passage="Market is actually $10B according to analysis"  # says $10B, not $25B
        )
        claim = builder.add_atomic_claim(
            statement="Market is $25B",  # claims $25B
            slide_number=4,
            visible_citation="Source: market_report.pdf p.14"
        )
        # Link $25B to evidence that says $10B
        for ns in claim.numeric_spans:
            if ns.value == "$25B":
                ns.evidence_id = eid
        
        verifier = AtomicContentVerifierV6()
        result = verifier.verify_atomic_claim(claim.numeric_spans[0], builder.evidences[0])
        
        st.json(result)
        if not result["verified"]:
            st.success("✓ V6 correctly FAILS - evidence says $10B but claim is $25B - adversarial test passed!")
            st.write(f"Reason: {result['reason']}")
        else:
            st.error("✗ V6 still passes - should fail!")

with col2:
    st.markdown("**Test 2: Tamper $25B hidden -> $999B visible**")
    if st.button("Run Test 2 - Should FAIL correctly"):
        from tamper_detector_v6 import TamperDetectorV6
        from pptx import Presentation
        from pptx.util import Inches
        
        builder = AtomicClaimBuilder()
        eid = builder.add_atomic_evidence(
            source_file="report.pdf",
            exact_location={"page": 4},
            exact_passage="Market is $25B"
        )
        claim = builder.add_atomic_claim(
            statement="Market is $25B",
            slide_number=2,
            visible_citation="Source: report.pdf p.4"
        )
        for ns in claim.numeric_spans:
            ns.evidence_id = eid
        
        # Create slide with tampered visible text $999B but hidden says $25B
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Market"
        slide.placeholders[1].text = "Market is $999B"  # Tampered visible!
        
        # Add hidden correct metadata
        from invisible_metadata_engine import InvisibleMetadataEngine
        inv_engine = InvisibleMetadataEngine()
        inv_engine.add_invisible_metadata_to_slide(slide, claim)  # Hidden says $25B
        
        pptx_path = "/tmp/test_tamper.pptx"
        prs.save(pptx_path)
        
        # Load and detect tamper
        prs_check = Presentation(pptx_path)
        detector = TamperDetectorV6()
        tamper_result = detector.detect_tampering(prs_check.slides[0], claim)
        
        st.json(tamper_result)
        if tamper_result["tamper_detected"]:
            st.success("✓ V6 correctly detects TAMPER - visible $999B != hidden $25B - export blocked!")
            for issue in tamper_result["issues"]:
                st.write(f"• {issue}")
        else:
            st.error("✗ V6 did not detect tamper!")

st.subheader("Full Pipeline Test - End-to-End")

uploaded = st.file_uploader("Upload PDF/Excel (real pipeline)", type=["pdf","xlsx"], accept_multiple_files=True)

if st.button("Run V6 Full Pipeline - Upload -> Extract -> Atomic Link -> Invisible -> QC -> Export Block if Fail", type="primary"):
    if not uploaded:
        st.warning("Upload files for full pipeline, or will use demo data")
        saved_paths = []
    else:
        saved_paths = []
        for f in uploaded:
            p = f"/tmp/{f.name}"
            with open(p, "wb") as out:
                out.write(f.getbuffer())
            saved_paths.append(p)
    
    engine = RealEngineV6FullPipeline(tavily_key=tavily_key, openai_key=openai_key)
    
    with st.status("V6 Full Pipeline - end-to-end", expanded=True) as status:
        st.write("Step 2: Extract with page/cell")
        facts = engine.step_2_extract(saved_paths) if saved_paths else [{"claim": "$25B", "source_file": "demo.pdf", "source_text": "Market $25B", "page_number": 14}]
        st.write(f"✓ {len(facts)} facts extracted, {len(engine.evidence_store)} evidences in store")
        
        st.write("Step 3: Tavily research (if key)")
        if tavily_key:
            research = engine.step_3_research("EV charging market")
            st.write(f"✓ Research done, total evidences: {len(engine.evidence_store)}")
        else:
            st.write("No Tavily key - using extracted evidences only")
        
        st.write("Step 5: Generate atomic claims with numeric spans")
        idea = "EV charging for apartments TAM $25B revenue $10M to $25M"
        slides, atomic_claims = engine.step_5_generate_atomic_claims(idea, facts)
        st.write(f"✓ {len(atomic_claims)} atomic claims, example: {atomic_claims[0].statement[:80] if atomic_claims else 'none'}")
        for ac in atomic_claims[:2]:
            for ns in ac.numeric_spans:
                st.caption(f"  {ns.value} [{ns.numeric_id}] -> {ns.evidence_id}")
        
        st.write("Step 7: Create deck with invisible metadata")
        template_path = "templates/minimal.pptx" if Path("templates/minimal.pptx").exists() else None
        output_path = "/tmp/citedeck_v6_full_pipeline.pptx"
        pptx_path = engine.step_7_create_deck_with_invisible_layer(template_path, slides, atomic_claims, output_path)
        st.write(f"✓ Deck created: {pptx_path} - visible clean, invisible atomic mapping in notes")
        
        st.write("Step 8: V6 QC - content verification + tamper detection + atomic")
        qc_result = engine.step_8_qc_v6_full(pptx_path)
        st.json(qc_result)
        
        if qc_result["can_publish"]:
            st.success(f"✓ {qc_result['trust_badge']} - Export allowed")
        else:
            st.error(f"✗ {qc_result['trust_badge']} - Export BLOCKED")
            for issue in qc_result["issues"]:
                st.write(f"• {issue}")
        
        status.update(label="V6 Pipeline done", state="complete")
    
    if qc_result.get("can_publish"):
        st.download_button("Download V6 Verified Deck (clean citations, auditable via notes)", open(pptx_path, "rb"), file_name="citedeck_v6_verified.pptx")
    else:
        st.warning("Export blocked due to verification failures - fix issues above")

st.markdown("---")
st.markdown("**V6 fixes your 4 issues + payment packaging:**")
st.markdown("- Content verification: checks passage actually contains claimed number, detects $25B vs $10B conflict")
st.markdown("- Tamper detection: compares visible $999B vs hidden $25B, blocks export")
st.markdown("- Production retrieval: scores evidence by exact value + keyword overlap + context, avoids false $10M mapping")
st.markdown("- Full pipeline: Upload -> Extract -> Research -> Atomic claims -> Invisible layer -> QC with tamper+content checks -> block if fail")
