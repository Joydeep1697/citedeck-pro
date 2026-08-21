from __future__ import annotations

from pathlib import Path
import tempfile
import unittest

from pptx import Presentation

from atomic_claim_evidence import AtomicClaimBuilder
from atomic_content_verifier_v6 import AtomicContentVerifierV6
from evidence_extractor_real import EvidenceExtractorReal
from evidence_retriever_v6 import EvidenceRetrieverV6
from real_engine_v6_full_pipeline import RealEngineV6FullPipeline


class ClaimDomainTests(unittest.TestCase):
    def setUp(self):
        self.builder = AtomicClaimBuilder()
        self.verifier = AtomicContentVerifierV6()

    def test_claim_without_evidence_stays_unmapped(self):
        claim = self.builder.add_atomic_claim("Revenue increased from $10M to $25M", 1, "Source: financials.xlsx")
        self.assertEqual([span.value for span in claim.numeric_spans], ["$10M", "$25M"])
        self.assertEqual(len(self.builder.evidences), 0)
        self.assertTrue(all(span.evidence_id is None for span in claim.numeric_spans))

    def test_synthetic_evidence_is_rejected(self):
        with self.assertRaises(ValueError):
            self.builder.add_atomic_evidence("pending_lookup", {"pending": True}, "Evidence for $25M")

    def test_number_matching_normalizes_units_without_substring_matches(self):
        self.assertTrue(self.verifier.does_passage_contain_number("$25B", "The market is worth $25 billion.")["supports"])
        self.assertFalse(self.verifier.does_passage_contain_number("$25B", "The market is worth $250B.")["supports"])
        self.assertFalse(self.verifier.does_passage_contain_number("25%", "The total value is 25.")["supports"])

    def test_retrieval_uses_exact_evidence_record_and_context(self):
        expense_id = self.builder.add_atomic_evidence("model.xlsx", {"cell": "B2", "cell_range": "Costs!B2"}, "Marketing expenses were $10M")
        revenue_id = self.builder.add_atomic_evidence("model.xlsx", {"cell": "B3", "cell_range": "Revenue!B3"}, "Subscription revenue was $10M")
        claim = self.builder.add_atomic_claim("Subscription revenue was $10M", 1, "Source: model.xlsx")
        EvidenceRetrieverV6().link_claim_atomically(claim, self.builder.evidences)
        self.assertNotEqual(expense_id, revenue_id)
        self.assertEqual(claim.numeric_spans[0].evidence_id, revenue_id)

    def test_unsupported_number_does_not_keep_old_mapping(self):
        self.builder.add_atomic_evidence("model.xlsx", {"cell": "B2"}, "Revenue was $10M")
        claim = self.builder.add_atomic_claim("Revenue was $25M", 1, "Source: model.xlsx")
        result = EvidenceRetrieverV6().link_claim_atomically(claim, self.builder.evidences)
        self.assertIsNone(claim.numeric_spans[0].evidence_id)
        self.assertFalse(result[0]["suitable"])


class StaticNarrator:
    def __init__(self, slides):
        self.slides = slides
        self.received_evidence = None

    def generate_defensible_deck(self, idea, facts, evidence):
        self.received_evidence = list(evidence)
        return self.slides


class PipelineTests(unittest.TestCase):
    def build_engine(self, slides=None):
        narrator = StaticNarrator(slides or [{"title": "Business performance", "bullets": ["Subscription revenue was $10M", "Customer growth reached 25%"], "citations": ["financials.xlsx"]}])
        engine = RealEngineV6FullPipeline(narrator=narrator, signing_key="unit-test-signing-secret")
        engine._add_fact_evidence({"claim": "$10M", "source_file": "financials.xlsx", "cell": "B2", "cell_range": "Revenue!B2", "source_text": "Subscription revenue was $10M"})
        engine._add_fact_evidence({"claim": "25%", "source_file": "financials.xlsx", "cell": "C2", "cell_range": "Growth!C2", "source_text": "Customer growth reached 25%"})
        return engine, narrator

    def create_deck(self, engine, path):
        slides, claims = engine.step_5_generate_atomic_claims("Create an investor presentation", [])
        engine.step_7_create_deck_with_invisible_layer(None, slides, claims, path)
        return slides, claims

    def test_all_bullets_are_embedded_verified_and_given_to_model(self):
        engine, narrator = self.build_engine()
        with tempfile.TemporaryDirectory() as directory:
            filename = Path(directory) / "deck.pptx"
            _, claims = self.create_deck(engine, filename)
            report = engine.step_8_qc_v6_full(filename)
            self.assertTrue(report["can_publish"], report["issues"])
            self.assertEqual(report["verified_claim_count"], 2)
            self.assertEqual(len(report["atomic_content_checks"]), 2)
            self.assertEqual(len(narrator.received_evidence), 2)
            slide = Presentation(filename).slides[0]
            notes = engine.invisible_engine.extract_invisible_metadata(slide)
            self.assertEqual(notes["citedeck_verification"]["claim_count"], len(claims))
            self.assertEqual(notes["citedeck_verification"]["integrity"]["algorithm"], "hmac-sha256")

    def test_visible_number_tampering_blocks_export(self):
        engine, _ = self.build_engine()
        with tempfile.TemporaryDirectory() as directory:
            filename = Path(directory) / "deck.pptx"
            self.create_deck(engine, filename)
            presentation = Presentation(filename)
            slide = presentation.slides[0]
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and "$10M" in shape.text:
                    shape.text = shape.text.replace("$10M", "$999M")
                    break
            presentation.save(filename)
            report = engine.step_8_qc_v6_full(filename)
            self.assertFalse(report["can_publish"])
            self.assertTrue(any("numbers" in issue or "missing or changed" in issue for issue in report["issues"]))

    def test_note_tampering_blocks_export(self):
        engine, _ = self.build_engine()
        with tempfile.TemporaryDirectory() as directory:
            filename = Path(directory) / "deck.pptx"
            self.create_deck(engine, filename)
            presentation = Presentation(filename)
            notes = presentation.slides[0].notes_slide.notes_text_frame
            notes.text = notes.text.replace("$10M", "$999M")
            presentation.save(filename)
            report = engine.step_8_qc_v6_full(filename)
            self.assertFalse(report["can_publish"])
            self.assertTrue(any("integrity" in issue for issue in report["issues"]))

    def test_unsupported_bullet_blocks_export(self):
        engine, _ = self.build_engine([{"title": "Unsupported", "bullets": ["Revenue was $999M"], "citations": ["financials.xlsx"]}])
        with tempfile.TemporaryDirectory() as directory:
            filename = Path(directory) / "deck.pptx"
            self.create_deck(engine, filename)
            report = engine.step_8_qc_v6_full(filename)
            self.assertFalse(report["can_publish"])
            self.assertTrue(any("independently sourced" in issue for issue in report["issues"]))

    def test_missing_evidence_fails_before_generation(self):
        engine = RealEngineV6FullPipeline(narrator=StaticNarrator([]))
        with self.assertRaises(ValueError):
            engine.step_5_generate_atomic_claims("Deck", [])

    def test_unaudited_title_number_is_rejected(self):
        engine, _ = self.build_engine([{"title": "Revenue in 2030", "bullets": ["Subscription revenue was $10M"], "citations": ["financials.xlsx"]}])
        with self.assertRaisesRegex(ValueError, "unaudited number"):
            engine.step_5_generate_atomic_claims("Deck", [])


class ExtractionTests(unittest.TestCase):
    def test_csv_rows_have_real_locations(self):
        with tempfile.TemporaryDirectory() as directory:
            filename = Path(directory) / "metrics.csv"
            filename.write_text("metric,amount\nrevenue,25000000\n", encoding="utf-8")
            facts = EvidenceExtractorReal().extract_csv_with_rows(filename)
            self.assertEqual(len(facts), 1)
            self.assertEqual(facts[0]["cell"], "R2C2")
            self.assertIn("25000000", facts[0]["source_text"])


if __name__ == "__main__":
    unittest.main()
