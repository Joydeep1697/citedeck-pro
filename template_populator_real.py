import os, json, re
from pptx import Presentation
from pptx.util import Inches, Pt

class TemplatePopulatorReal:
    """Adaptive template populator - works even if placeholders missing"""
    
    def find_placeholder(self, slide, placeholder_names):
        """Search for {{title}}, {{bullets}}, {{citation}} or fallback to layout"""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            for name in placeholder_names:
                if name in text or name.strip("{}") in text.lower():
                    return shape
        return None
    
    def populate_slide_adaptive(self, slide, title, bullets, citations, chart_path=None):
        """Populate with fallbacks - creates textboxes if placeholders not found"""
        
        # 1. Title
        title_shape = self.find_placeholder(slide, ["{{title}}", "title", "Title"])
        if title_shape:
            title_shape.text_frame.text = title
            title_shape.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
            title_shape.text_frame.paragraphs[0].runs[0].font.bold = True
        else:
            # Create title if not found - check if slide has title placeholder
            if slide.shapes.title:
                slide.shapes.title.text = title
            else:
                # Create new textbox for title
                textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                tf = textbox.text_frame
                tf.text = title
                tf.paragraphs[0].runs[0].font.size = Pt(28)
                tf.paragraphs[0].runs[0].font.bold = True
        
        # 2. Bullets
        bullets_shape = self.find_placeholder(slide, ["{{bullets}}", "{{content}}", "bullets", "content"])
        bullets_text = "\n".join([f"• {b}" for b in bullets])
        
        if bullets_shape:
            bullets_shape.text_frame.text = bullets_text
        else:
            # Find content placeholder or create
            content_added = False
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title and "Source:" not in shape.text:
                    # Use first non-title text box as content
                    if len(shape.text_frame.text) < 50:  # likely placeholder
                        shape.text_frame.text = bullets_text
                        content_added = True
                        break
            
            if not content_added:
                textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.5), Inches(3.5))
                tf = textbox.text_frame
                tf.text = bullets_text
                tf.word_wrap = True
        
        # 3. Citations - ALWAYS add citation footer (critical for trust)
        citation_text = f"Sources: {', '.join(citations[:3])}" if citations else "Source: Analysis"
        citation_shape = self.find_placeholder(slide, ["{{citation}}", "{{sources}}", "citation", "source"])
        
        if citation_shape:
            citation_shape.text_frame.text = citation_text
            citation_shape.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
        else:
            # Always create citation footer
            footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.4))
            tf = footer.text_frame
            tf.text = citation_text
            tf.paragraphs[0].runs[0].font.size = Pt(9)
            tf.paragraphs[0].runs[0].font.italic = True
        
        # 4. Chart if provided
        if chart_path and os.path.exists(chart_path):
            try:
                slide.shapes.add_picture(chart_path, Inches(6.2), Inches(1.2), Inches(3.5), Inches(3.5))
            except:
                pass

    def populate_presentation(self, template_path, slides_data, charts, output_path):
        """Real population with adaptive layout"""
        if not os.path.exists(template_path):
            # Fallback to minimal template creation
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
        else:
            prs = Presentation(template_path)
        
        # Ensure enough slides
        while len(prs.slides) < len(slides_data):
            # Use blank or title+content layout
            try:
                layout = prs.slide_layouts[1]  # title + content
            except:
                layout = prs.slide_layouts[0]
            prs.slides.add_slide(layout)
        
        for i, slide_data in enumerate(slides_data):
            if i >= len(prs.slides):
                break
            slide = prs.slides[i]
            
            chart_path = None
            if charts and i < len(charts) and "chart_path" in charts[i]:
                chart_path = charts[i]["chart_path"]
            
            self.populate_slide_adaptive(
                slide,
                slide_data.get("title", f"Slide {i+1}"),
                slide_data.get("bullets", []),
                slide_data.get("citations", []),
                chart_path
            )
        
        prs.save(output_path)
        return output_path
