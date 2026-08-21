"""
V5 Invisible Metadata Layer

Fixes: Embedded IDs should not be visible in final deck

Before V4: Slide shows "Market $25B [CLAIM-042] [EVIDENCE-117]" visible to customer - unprofessional

After V5:
  Visible slide: "Market will reach $25B by 2030" + "Source: Gartner, 2026, p.14"
  Hidden in notes: JSON with CLAIM-042 -> EVIDENCE-117 mapping for QC/audit
  
User sees clean citation, CiteDeck internally maintains verification layer
"""

import json
from pptx import Presentation
from pptx.util import Inches

class InvisibleMetadataEngine:
    """
    Stores CLAIM->EVIDENCE mapping in invisible layer:
    - Speaker notes (primary - works in all PPTX viewers, hidden from audience)
    - Custom XML properties (secondary - for machine verification)
    """
    
    def add_invisible_metadata_to_slide(self, slide, atomic_claim):
        """
        Visible: clean citation
        Invisible: atomic mapping in notes
        """
        # 1. Visible layer - what customer sees
        visible_citation = atomic_claim.visible_citation  # "Source: Gartner, 2026, p.14"
        
        # Find or create citation shape for visible citation
        citation_added = False
        for shape in slide.shapes:
            if shape.has_text_frame and "Source:" in shape.text_frame.text:
                shape.text_frame.text = visible_citation
                citation_added = True
                break
        
        if not citation_added:
            # Add clean visible citation at bottom
            footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.4))
            footer.text_frame.text = visible_citation
            footer.text_frame.paragraphs[0].runs[0].font.size = Inches(0.12)  # small
        
        # 2. Invisible layer - speaker notes (hidden from audience, visible in presenter view)
        try:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            
            # Build invisible audit JSON
            invisible_audit = {
                "citedeck_verification": {
                    "claim_id": atomic_claim.claim_id,
                    "statement": atomic_claim.statement,
                    "atomic_numeric_mapping": [
                        {
                            "numeric_id": ns.numeric_id,
                            "value": ns.value,
                            "char_span_in_claim": [ns.char_start, ns.char_end],
                            "evidence_id": ns.evidence_id,
                            "verification": ns.verification_status
                        }
                        for ns in atomic_claim.numeric_spans
                    ],
                    "evidence_ids": atomic_claim.evidence_ids,
                    "visible_citation": atomic_claim.visible_citation,
                    "audit_trail": "This slide is auditable - each number maps to exact source location"
                }
            }
            
            # Store in notes - invisible during presentation, visible for audit
            notes_text_frame.text = json.dumps(invisible_audit, indent=2)
        except Exception as e:
            print(f"Notes slide not available, using alternative: {e}")
            # Fallback: store in slide tags (custom XML)
            try:
                # Use slide's custom properties
                slide.tags.add("citedeck_claim_id", atomic_claim.claim_id)
                slide.tags.add("citedeck_evidence_ids", ",".join(atomic_claim.evidence_ids))
                slide.tags.add("citedeck_atomic_mapping", json.dumps([{"value": ns.value, "evidence": ns.evidence_id} for ns in atomic_claim.numeric_spans]))
            except:
                pass
    
    def extract_invisible_metadata(self, slide):
        """QC reads invisible layer, not visible text"""
        try:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text
            if "citedeck_verification" in notes_text:
                return json.loads(notes_text)
        except:
            pass
        
        # Try tags fallback
        try:
            if hasattr(slide, 'tags'):
                claim_id = slide.tags.get("citedeck_claim_id")
                if claim_id:
                    return {"citedeck_verification": {"claim_id": claim_id}}
        except:
            pass
        
        return None
    
    def create_professional_slide(self, prs, title, statement, visible_citation, atomic_claim):
        """Creates slide that looks professional - no visible IDs"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = statement  # Clean, no IDs
        
        # Add visible clean citation
        footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.4))
        footer.text_frame.text = visible_citation  # "Source: Gartner, 2026, p.14" - clean
        
        # Add invisible audit layer
        self.add_invisible_metadata_to_slide(slide, atomic_claim)
        
        return slide

if __name__ == "__main__":
    from atomic_claim_evidence import AtomicClaimBuilder
    
    builder = AtomicClaimBuilder()
    eid = builder.add_atomic_evidence(
        source_file="market_report.pdf",
        source_type="pdf",
        exact_location={"page": 14},
        exact_passage="Market will reach $25B"
    )
    claim = builder.add_atomic_claim(
        statement="Market will reach $25B by 2030",
        slide_number=7,
        visible_citation="Source: Gartner, 2026, p.14"
    )
    builder.link_number_to_evidence(claim.claim_id, "$25B", eid)
    
    prs = Presentation()
    engine = InvisibleMetadataEngine()
    engine.create_professional_slide(prs, "Market", claim.statement, claim.visible_citation, claim)
    prs.save("/tmp/test_invisible_metadata.pptx")
    print("Created deck with invisible metadata - visible clean, hidden auditable")
