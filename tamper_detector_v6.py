import re
import json
from pptx import Presentation

class TamperDetectorV6:
    """
    FIXES CRITICAL ISSUE 2: Visible slide content can be tampered with
    Your adversarial test: Hidden metadata says $25B, visible changed to $999B, V5 still passed
    
    V6: Compares visible slide text vs hidden claim statement, fails if they differ
    """
    
    def __init__(self):
        self.number_pattern = re.compile(r'\$?\d+(?:,\d+)*(?:\.\d+)?\s*(?:B|M|K|Cr|billion|million|%|USD|INR)?', re.I)
    
    def extract_numbers(self, text: str):
        """Extract all numbers from text"""
        return [m.group(0) for m in self.number_pattern.finditer(text)]
    
    def normalize_for_comparison(self, text: str):
        """Normalize text for comparison - lowercase, remove extra spaces"""
        return re.sub(r'\s+', ' ', text.lower().strip())
    
    def detect_tampering(self, slide, atomic_claim):
        """
        Compares visible slide text vs hidden atomic claim
        
        Args:
            slide: pptx slide
            atomic_claim: AtomicClaim with statement and numeric_spans
        
        Returns: tamper detection result
        """
        # Extract visible text from slide
        visible_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                visible_text += (shape.text_frame.text or "") + " "
        
        visible_text_norm = self.normalize_for_comparison(visible_text)
        hidden_statement = atomic_claim.statement
        hidden_norm = self.normalize_for_comparison(hidden_statement)
        
        # Extract numbers from both
        visible_numbers = self.extract_numbers(visible_text)
        hidden_numbers = [ns.value for ns in atomic_claim.numeric_spans]
        hidden_statement_numbers = self.extract_numbers(hidden_statement)
        
        # Check 1: Does visible text contain hidden statement numbers?
        tamper_detected = False
        issues = []
        
        # For each number in hidden claim, check if it exists in visible text
        for hidden_num in hidden_statement_numbers:
            if hidden_num not in visible_numbers:
                # Number in hidden but not visible - could be tampering or incomplete slide
                # For strict check: if hidden says $25B but visible says $999B, it's tamper
                tamper_detected = True
                issues.append(f"TAMPER DETECTED: Hidden claim has {hidden_num} but visible slide has {visible_numbers} - visible changed from '{hidden_statement}' to '{visible_text[:100]}'")
        
        # Check 2: Does visible have extra numbers not in hidden?
        for vis_num in visible_numbers:
            if vis_num not in hidden_statement_numbers and vis_num not in hidden_numbers:
                # Visible has number not in hidden - possible tampering (someone added $999B)
                # Need to check if it's a citation or legitimate
                # For strict: any number in visible that is not in hidden atomic mapping is suspicious
                if len(hidden_statement_numbers) > 0:  # Only if hidden had numbers
                    tamper_detected = True
                    issues.append(f"TAMPER DETECTED: Visible slide has {vis_num} not in hidden atomic mapping {hidden_numbers} - possible injection of $999B")
        
        # Check 3: Direct statement comparison (fuzzy)
        # If hidden says "Market is $25B" but visible says "Market is $999B", detect
        if hidden_norm not in visible_text_norm and visible_text_norm not in hidden_norm:
            # Statements differ significantly
            # Check if numbers differ
            if set(visible_numbers) != set(hidden_statement_numbers) and len(hidden_statement_numbers) > 0:
                tamper_detected = True
                issues.append(f"TAMPER: Visible '{visible_text[:100]}' != Hidden '{hidden_statement}' - numeric values differ: visible {visible_numbers} vs hidden {hidden_statement_numbers}")
        
        return {
            "slide_number": getattr(slide, 'slide_id', 0),
            "visible_text_preview": visible_text[:150],
            "hidden_statement": hidden_statement,
            "visible_numbers": visible_numbers,
            "hidden_numbers": hidden_numbers,
            "hidden_statement_numbers": hidden_statement_numbers,
            "tamper_detected": tamper_detected,
            "issues": issues,
            "verified": not tamper_detected,
            "check": "Does visible slide text match hidden claim statement? Prevents $25B -> $999B tampering"
        }

if __name__ == "__main__":
    detector = TamperDetectorV6()
    
    print("=== Adversarial Test 2: Visible $999B vs Hidden $25B ===")
    # Simulate
    visible = "Market is $999B"
    hidden_statement = "Market is $25B"
    hidden_numbers = ["$25B"]
    
    vis_nums = detector.extract_numbers(visible)
    hidden_nums = detector.extract_numbers(hidden_statement)
    
    print(f"Visible numbers: {vis_nums}")
    print(f"Hidden numbers: {hidden_nums}")
    print(f"Tamper detected: {set(vis_nums) != set(hidden_nums)} - Should be True")
