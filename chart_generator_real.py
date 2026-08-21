import os, json, re
from pathlib import Path

class ChartGeneratorReal:
    """Real chart generation: dataframe -> matplotlib -> png -> insert into PPTX"""
    
    def __init__(self, output_dir="/tmp/citedeck_charts"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
    
    def generate_charts_from_excel(self, excel_path, max_charts=3):
        """Creates real PNG charts from Excel"""
        try:
            import pandas as pd
            import matplotlib
            matplotlib.use('Agg')
            import matplotlib.pyplot as plt
            
            xl = pd.ExcelFile(excel_path)
            charts = []
            
            for sheet in xl.sheet_names[:2]:
                df = xl.parse(sheet, nrows=30)
                if df.shape[0] < 2 or df.shape[1] < 2:
                    continue
                
                # Find numeric columns
                numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
                if not numeric_cols:
                    continue
                
                # Try to find label column (first non-numeric)
                label_col = None
                for c in df.columns:
                    if c not in numeric_cols:
                        label_col = c
                        break
                
                for num_col in numeric_cols[:2]:
                    plt.figure(figsize=(6,4))
                    if label_col:
                        plt.bar(df[label_col].astype(str)[:10], df[num_col][:10])
                        plt.xticks(rotation=30, ha='right')
                    else:
                        plt.plot(df[num_col][:15], marker='o')
                    plt.title(f"{num_col} from {sheet}")
                    plt.tight_layout()
                    
                    chart_path = self.output_dir / f"{Path(excel_path).stem}_{sheet}_{num_col}.png"
                    plt.savefig(chart_path, dpi=150)
                    plt.close()
                    
                    charts.append({
                        "chart_path": str(chart_path),
                        "source_file": os.path.basename(excel_path),
                        "sheet": sheet,
                        "column": num_col,
                        "cell_range": f"{sheet}!{num_col}",
                        "type": "bar" if label_col else "line",
                        "verification": f"Chart from {sheet}!{num_col} - {chart_path}"
                    })
                    if len(charts) >= max_charts:
                        break
                if len(charts) >= max_charts:
                    break
            
            return charts
        except Exception as e:
            return [{"error": str(e), "source_file": os.path.basename(excel_path)}]

    def insert_chart_into_pptx(self, pptx_path, chart_info, slide_index=5):
        """Actually inserts PNG into PPTX slide"""
        try:
            from pptx import Presentation
            from pptx.util import Inches
            
            prs = Presentation(pptx_path)
            if slide_index >= len(prs.slides):
                slide_index = len(prs.slides) - 1
            
            slide = prs.slides[slide_index]
            # Add picture
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(4)
            
            slide.shapes.add_picture(chart_info["chart_path"], left, top, width=width, height=height)
            
            # Add citation textbox
            textbox = slide.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(9), Inches(0.5))
            tf = textbox.text_frame
            tf.text = f"Source: {chart_info['source_file']} | {chart_info.get('cell_range','')} | Chart generated from your file"
            
            prs.save(pptx_path)
            return True
        except Exception as e:
            print(f"Insert chart error: {e}")
            return False
